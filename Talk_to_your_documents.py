import io
import os
import streamlit as st 
from langchain.chat_models import AzureChatOpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import FAISS
import pandas as pd
from langchain.agents import create_csv_agent
#from langchain_experimental.agents import create_csv_agent
from PyPDF2 import PdfReader
from pptx import Presentation
from azure.identity import ClientSecretCredential
from azure.keyvault.secrets import SecretClient


# Customize the layout
st.title("Document QnA")

uploadType = st.radio(
    "What are you uploading...",
    ('CSV', 'TEXT', 'PDF', 'PPT'
    ), horizontal=True)

TENANT= "c0bd3319-c6f8-48e9-95ff-112a47c538b7"
CLIENT_ID = "188c2ca6-030b-4ed4-9b62-52181e0369ba"
CLIENT_SECRET= st.secrets.CLIENT_SECRET
credential = ClientSecretCredential(TENANT,CLIENT_ID,CLIENT_SECRET)
VAULT_URL= "https://tsa-dev-akv.vault.azure.net/"
client = SecretClient(vault_url=VAULT_URL, credential=credential)
openai_key = client.get_secret("GenAIBIMInternalCapabilityOpenAIKey")

os.environ["OPENAI_API_TYPE"] = "azure"
os.environ["OPENAI_API_BASE"] = "https://biminternalcapability.openai.azure.com/"
os.environ["OPENAI_API_KEY"] = openai_key.value
os.environ["OPENAI_API_VERSION"] = "2023-07-01-preview"

# function for writing uploaded file in temp
def write_text_file(content, file_path):
    try:
        with open(file_path, 'w') as file:
            file.write(content)
        return True
    except Exception as e:
        print(f"Error occurred while writing the file: {e}")
        return False

# set prompt template
prompt_template = """Use the following pieces of context to answer the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.

{context}

Question: {question}
Answer:"""
prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])

# initialize the LLM & Embeddings
llm = AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4")
embeddings = OpenAIEmbeddings(deployment="embeddings", chunk_size=1)
llm_chain = LLMChain(llm=llm, prompt=prompt)

if uploadType == 'CSV':
    uploaded_file = st.file_uploader("Upload the file", type="csv")
elif uploadType == 'TEXT':
    uploaded_file = st.file_uploader("Upload the file", type="txt")
elif uploadType == 'PPT':
    uploaded_file = st.file_uploader("Upload the file", type="pptx")
else:
    uploaded_file = st.file_uploader("Upload the file", type="pdf")

#for TEXT files
if uploaded_file is not None and uploadType == 'TEXT':
    content = uploaded_file.read().decode('utf-8')
    # st.write(content)
    file_path = "./temp/file.txt"
    write_text_file(content, file_path)   
    
    #loader = TextLoader(file_path)
    #docs = loader.load()    

    # Open the text file and read the text.
    loader = open(file_path, "r")
    docs = loader.read()
    #text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    
    text_splitter = CharacterTextSplitter(        
    separator = "\n",
    chunk_size = 1000,
    chunk_overlap  = 200,
    length_function = len,
    )
    
    texts = text_splitter.split_text(docs)
    docSearch = FAISS.from_texts(texts, embeddings)
    #chain = load_qa_chain(AzureOpenAI(deployment_name="gpt-35-turbo-dev", model_name="gpt-35-turbo"), chain_type="stuff")
    st.success("File Loaded Successfully!!")
    
    # Query through LLM    
    question = st.text_input("Ask something from the file", placeholder="Find something similar in the file...", disabled=not uploaded_file,)    
    print(question)
    if question:
        similar_doc = docSearch.similarity_search(question, k=1)
        context = similar_doc[0].page_content
        query_llm = LLMChain(llm=llm, prompt=prompt)
        response = query_llm.run({"context": context, "question": question})        
        #response = chain.run(input_documents=similar_doc[0].page_content, question=question)
        st.write(response)

#for CSV files
if uploaded_file is not None and uploadType == 'CSV':
    df_csv = pd.read_csv(uploaded_file, encoding= 'unicode_escape')
    df_csv_top10 = df_csv.head(2)
    st.write("Sample records displayed...")
    st.dataframe(df_csv_top10)
    
    file_path = "./temp/file.csv"
    df_csv.to_csv(file_path)
    st.success("File Loaded Successfully!!")

    #agent = create_csv_agent(AzureOpenAI(deployment_name="text-davinci-003-dev", model_name="text-davinci-003"), file_path)
    agent = create_csv_agent(AzureChatOpenAI(deployment_name="gpt-4", model_name="gpt-4"), file_path)

    # Query through LLM    
    question = st.text_input("Ask something from the file", placeholder="Find something similar in the file...", disabled=not uploaded_file,)    
    print(question)
    if question:
        response = agent.run(question)
        st.write(response)

#for PDF files
if uploaded_file is not None and uploadType == 'PDF':
    pdf_read = PdfReader(uploaded_file)
    #file_path = "./temp/file.csv"
    
    raw_text = ''
    for i, page in enumerate(pdf_read.pages):
        text = page.extract_text()
        if text:
            raw_text += text
        
    text_splitter = CharacterTextSplitter(        
    separator = "\n",
    chunk_size = 1000,
    chunk_overlap  = 200,
    length_function = len,
    )
    
    pdf_texts = text_splitter.split_text(raw_text)
    pdfdocSearch = FAISS.from_texts(pdf_texts, embeddings)
    #chain = load_qa_chain(AzureOpenAI(deployment_name="gpt-35-turbo-dev", model_name="gpt-35-turbo"), chain_type="stuff")
    st.success("File Loaded Successfully!!")
    
    # Query through LLM    
    question = st.text_input("Ask something from the file", placeholder="Find something similar in the file...", disabled=not uploaded_file,)    
    print(question)
    if question:
        similar_doc = pdfdocSearch.similarity_search(question, k=1)
        context = similar_doc[0].page_content
        query_llm = LLMChain(llm=llm, prompt=prompt)
        response = query_llm.run({"context": context, "question": question})        
        #response = chain.run(input_documents=similar_doc[0].page_content, question=question)
        st.write(response)

#for PPT files
if uploaded_file is not None and uploadType == 'PPT':
    ppt_content = uploaded_file.read()
    presentation = Presentation(io.BytesIO(ppt_content))
    ppt_text = ""

    text_splitter = CharacterTextSplitter(        
    separator = "\n",
    chunk_size = 1000,
    chunk_overlap  = 200,
    length_function = len,
    )

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += shape.text

    # Split the text into chunks for processing
    ppt_texts = text_splitter.split_text(ppt_text)

    # Create FAISS index
    ppt_docSearch = FAISS.from_texts(ppt_texts, embeddings)
    st.success("File Loaded Successfully!!")

    # Query through LLM
    question = st.text_input("Ask something from the file", placeholder="Find something similar in the file...", disabled=not uploaded_file,)    
    print(question)
    if question:
        similar_doc = ppt_docSearch.similarity_search(question, k=1)
        context = similar_doc[0].page_content
        query_llm = LLMChain(llm=llm, prompt=prompt)
        response = query_llm.run({"context": context, "question": question})
        st.write(response)